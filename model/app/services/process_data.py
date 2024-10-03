from langchain_community.chat_models import ChatOpenAI 
from langchain_core.prompts import PromptTemplate, ChatPromptTemplate 
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from llama_parse import LlamaParse
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain_community.utilities import GoogleSerperAPIWrapper
from langchain.tools import WikipediaQueryRun
from langchain_community.utilities import WikipediaAPIWrapper
from langchain.tools.retriever import create_retriever_tool
from langchain.agents import AgentExecutor
from langchain.agents import create_openai_tools_agent
from pydantic import BaseModel
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain.tools import Tool
from langchain.tools.retriever import create_retriever_tool
from langchain_core.tools import StructuredTool
import fitz  

from dotenv import load_dotenv
import os
load_dotenv(dotenv_path='.env', override=True)
llama_api_key = os.getenv("LLAMA_CLOUD_API_KEY")
openai_api_key = os.getenv("OPENAI_API_KEY")

import nest_asyncio
nest_asyncio.apply()

from pptx import Presentation
def extract_text_from_ppt(ppt_file):
    """
    Extracts text from a PowerPoint file (.pptx)
    
    :param ppt_file: The path to the PowerPoint fil
    :return: The text extracted from the PowerPoint file
    """
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_pdf_text(pdf_path: str) -> str:
    """Extract text from a PDF file."""
    pdf_text = ""
    with fitz.open(pdf_path) as pdf_document:
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num) 
            pdf_text += page.get_text() 
    return pdf_text

from langchain.schema import Document
import os
agent_executor = None

from app.prompt import set_agent_executor

def setup_database() -> None:

    global agent_executor
    """Setup the database of documents."""
    pdf_folder = "./data/pdf_files/"
    ppt_folder = "./data/ppt_files/"

    # Load all PDFs from the folder
    pdf_documents = []
    for pdf_file in os.listdir(pdf_folder):
        if pdf_file.endswith(".pdf"):
            pdf_path = os.path.join(pdf_folder, pdf_file)
            pdf_content = extract_pdf_text(pdf_path)
            pdf_documents.append(Document(page_content=pdf_content))

    # Load all PPTs from the folder
    ppt_documents = []
    for ppt_file in os.listdir(ppt_folder):
        if ppt_file.endswith(".pptx"):
            ppt_path = os.path.join(ppt_folder, ppt_file)
            ppt_text = extract_text_from_ppt(ppt_path)
            ppt_documents.append(Document(page_content=ppt_text))

    # Split the PDFs and PPTs into chunks
    pdf_chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(pdf_documents)
    ppt_chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(ppt_documents)

    pdfvectorstore = Chroma.from_documents(documents=pdf_chunks, embedding=OpenAIEmbeddings())
    pptvectorstore = Chroma.from_documents(documents=ppt_chunks, embedding=OpenAIEmbeddings())

    pdfretriever = pdfvectorstore.as_retriever()
    pptretriever = pptvectorstore.as_retriever()

    search = GoogleSerperAPIWrapper()
    wikipedia = WikipediaAPIWrapper()

    web_search_tool = StructuredTool.from_function(name="web_search", func=search.run, description="Search the web for current information. Use this as a last resort to fill gaps in knowledge not covered by local documents or Wikipedia.")
    wikipedia_tool = WikipediaQueryRun(api_wrapper=wikipedia, description="Search Wikipedia for additional information if local documents don't provide sufficient details.")
    pptretriever_tool = create_retriever_tool(pptretriever, "ppt_search", description="Search for information in local PowerPoint presentations. Always use this in conjunction with pdf_search for comprehensive local information retrieval.")
    pdfretriever_tool = create_retriever_tool(pdfretriever, "pdf_search", description="Search for information in local PDF documents. Always use this in conjunction with ppt_search for comprehensive local information retrieval.")
    
    tools=[pdfretriever_tool,pptretriever_tool, wikipedia_tool, web_search_tool]
    llm = ChatOpenAI(model="gpt-3.5-turbo-0125", temperature=0)
    main_prompt=ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "You are a highly capable assistant with access to various information sources. "
                "Your primary goal is to provide comprehensive and accurate responses by combining information from all available sources. "
                "If it is a general greeting then greet properly without going into the tools"
                "When answering queries or completing tasks:\n"
                "1. Always start by searching both PowerPoint presentations and PDF documents simultaneously.\n"
                "2. Combine and synthesize information from both local sources to form a complete answer.\n"
                "3. If the local sources don't provide sufficient information, consult Wikipedia.\n"
                "4. Use web search as a last resort to fill any remaining gaps in information.\n"
                "5. Clearly indicate which sources you've used in your response.\n"
                "6. If asked to create content (like quizzes or summaries), use information from all relevant sources.\n"
                "Remember, your strength lies in your ability to integrate information from multiple sources effectively.",
            ),
            ("user", "{input}"),
            MessagesPlaceholder(variable_name="agent_scratchpad"),
        ]
    )
    
    main_prompt = ChatPromptTemplate.from_messages([
        MessagesPlaceholder(variable_name="agent_scratchpad"),
        ("system", "Your assistant configuration."),
        ("user", "{input}"),
    ])

    agent = create_openai_tools_agent(llm, tools, main_prompt)
    agent_exec=AgentExecutor(agent=agent,tools=tools,verbose=True)
    set_agent_executor(agent_exec)
    print("Agent Executor initialized.")

__all__ = ['setup_database', 'agent_executor']