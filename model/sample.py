from dotenv import load_dotenv
import os
load_dotenv(dotenv_path='.env', override=True)

from langchain.chat_models import ChatOpenAI
from langchain import PromptTemplate
from langchain.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from llama_parse import LlamaParse
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain.agents import create_openai_tools_agent
from langchain.schema import Document
import os
from pptx import Presentation

import nest_asyncio
nest_asyncio.apply()

def extract_text_from_ppt(ppt_file):
    """
    Extracts text from a PowerPoint file (.pptx)
    
    :param ppt_file: The path to the PowerPoint file
    :return: The text extracted from the PowerPoint file
    """
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text


def set_up_database():
    pdf_folder = "./data/pdf_files/"
    ppt_folder = "./data/ppt_files/"

    pdf_documents = []
    ppt_documents = []

    # Load all PDFs from the folder
    for pdf_file in os.listdir(pdf_folder):
        if pdf_file.endswith(".pdf"):
            pdf_path = os.path.join(pdf_folder, pdf_file)
            pdf_content = LlamaParse(result_type="markdown").load_data(pdf_path)
            pdf_documents.append(Document(page_content=pdf_content[0].text))

    # Load all PPTs from the folder
    for ppt_file in os.listdir(ppt_folder):
        if ppt_file.endswith(".pptx"):
            ppt_path = os.path.join(ppt_folder, ppt_file)
            ppt_text = extract_text_from_ppt(ppt_path)
            ppt_documents.append(Document(page_content=ppt_text))

    pdf_chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(pdf_documents)
    ppt_chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(ppt_documents)

    pdfvectorstore = Chroma.from_documents(documents=pdf_chunks, embedding=OpenAIEmbeddings())
    pptvectorstore = Chroma.from_documents(documents=ppt_chunks, embedding=OpenAIEmbeddings())

    pdfretriever = pdfvectorstore.as_retriever()
    pptretriever = pptvectorstore.as_retriever()

    from langchain_community.utilities import GoogleSerperAPIWrapper
    from langchain.tools import Tool
    search = GoogleSerperAPIWrapper()

    from langchain.tools import WikipediaQueryRun
    from langchain_community.utilities import WikipediaAPIWrapper

    wikipedia = WikipediaAPIWrapper()

    from langchain.tools.retriever import create_retriever_tool
    from langchain.utilities import WikipediaAPIWrapper

    web_search_tool = Tool(name="Web Search", func=search.run, description="Useful for searching the web for current information on various topics")
    wikipedia_tool = WikipediaQueryRun(api_wrapper=wikipedia)
    pptretriever_tool = create_retriever_tool(pptretriever, "ppt_search", description="PPT search")
    pdfretriever_tool = create_retriever_tool(pdfretriever, "pdf_search", description="PDF search")

    from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
    tools=[pdfretriever_tool,pptretriever_tool,web_search_tool,wikipedia_tool]
    llm = ChatOpenAI(model="gpt-3.5-turbo-0125", temperature=0)
    main_prompt=ChatPromptTemplate.from_messages(
        [
            (
                "system",
                "You are a very powerful assistant with access to various information sources. "
                "Always prioritize searching local documents before using external sources. "
                "Follow this order strictly when searching for information:\n"
                "1. Search PowerPoint presentations (use ppt_search tool)\n"
                "2. Search PDF documents (use pdf_search tool)\n"
                "3. Check Wikipedia (use wikipedia_search tool)\n"
                "4. Perform a web search (use web_search tool) as a last resort\n"
                "Only move to the next source if the information is not found in the current one. "
                "Always explain which source you're using and why."
            ),
            ("user", "{input}"),
            MessagesPlaceholder(variable_name="agent_scratchpad"),
        ]
    )


    agent=create_openai_tools_agent(llm,tools,main_prompt)

    from langchain.agents import AgentExecutor
    agent_executor=AgentExecutor(agent=agent,tools=tools,verbose=True)